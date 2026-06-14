"""
Locate, read and persist generated documents.

Every deliverable is written to GENERATED_DIR. Alongside it we keep a *source
sidecar* (the original Markdown or Python the document was built from) under a
hidden ``.sources`` subfolder, so a document can be re-read and modified later
without lossy text extraction from the binary. The subfolder is hidden so it
never shows up in the downloadable-files listing.
"""
from __future__ import annotations

import re
from pathlib import Path
from datetime import datetime, timezone

from app.config import GENERATED_DIR
from app.logging_config import logger

_SOURCES_DIR = GENERATED_DIR / ".sources"


def build_dest(formato: str, nombre: str) -> Path:
    """Return a unique, sanitized destination path inside GENERATED_DIR."""
    fmt = formato.lower().strip()
    safe_name = re.sub(r"[^\w\-]", "_", nombre) or "document"
    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    return GENERATED_DIR / f"{safe_name}_{timestamp}.{fmt}"


def _sources_dir() -> Path:
    _SOURCES_DIR.mkdir(parents=True, exist_ok=True)
    return _SOURCES_DIR


def save_source_sidecar(filename: str, source: str) -> None:
    """Persist the source (Markdown/Python) a document was generated from."""
    try:
        (_sources_dir() / f"{filename}.source.md").write_text(source, encoding="utf-8")
    except Exception:  # pragma: no cover - sidecar is best-effort
        logger.warning("Could not save the source sidecar for %s", filename)


def _resolve_generated(nombre: str) -> Path | None:
    """Locate a generated file by name (with or without extension).

    If several files match (same name, different timestamp) the most recent one
    is returned.
    """
    base = re.sub(r"\.source\.md$", "", nombre.strip())
    base = Path(base).name  # never accept paths
    candidates: list[Path] = []
    if (GENERATED_DIR / base).is_file():
        candidates.append(GENERATED_DIR / base)
    stem = re.sub(r"\.[^.]+$", "", base)
    for p in GENERATED_DIR.glob(f"{stem}*"):
        if p.is_file() and not p.name.endswith(".source.md"):
            candidates.append(p)
    if not candidates:
        return None
    return max(set(candidates), key=lambda p: p.stat().st_mtime)


def read_generated(nombre: str) -> str:
    """Return the text content of an already-generated document.

    Prefers the source sidecar (faithful re-read of what produced the file); if
    absent, extracts text from the binary according to its format.
    """
    target = _resolve_generated(nombre)
    if target is None:
        return (
            f"No generated document named '{nombre}' was found. "
            "Check the exact file name."
        )

    sidecar = _sources_dir() / f"{target.name}.source.md"
    if sidecar.is_file():
        return sidecar.read_text(encoding="utf-8")

    fmt = target.suffix.lower().lstrip(".")
    try:
        if fmt in ("md", "txt", "html"):
            return target.read_text(encoding="utf-8")

        if fmt == "pdf":
            import fitz  # PyMuPDF

            doc = fitz.open(str(target))
            return "\n".join(page.get_text() for page in doc)

        if fmt == "docx":
            from docx import Document

            doc = Document(str(target))
            return "\n".join(p.text for p in doc.paragraphs)

        if fmt == "xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(str(target), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"# {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    lines.append("\t".join("" if c is None else str(c) for c in row))
            return "\n".join(lines)

        if fmt == "pptx":
            from pptx import Presentation

            prs = Presentation(str(target))
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"## Slide {i}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        lines.append(shape.text_frame.text)
            return "\n".join(lines)

    except Exception as exc:  # pragma: no cover - extraction is best-effort
        logger.exception("Error reading generated document '%s': %s", target.name, exc)
        return (
            f"Could not extract the text from '{target.name}' (format {fmt}): {exc}. "
            "If you need to modify it, regenerate it from the source content."
        )

    return target.read_text(encoding="utf-8", errors="replace")
